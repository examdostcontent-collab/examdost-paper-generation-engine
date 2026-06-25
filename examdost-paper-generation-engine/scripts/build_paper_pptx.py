"""
build_paper_pptx.py — render a generated paper (paper.json) as a branded PPTX
question bank (one question per slide).

Usage:
    python build_paper_pptx.py <paper.json> <output.pptx>

Honours meta.layout ("combined" | "split"), meta.metadata_tags, meta.show_solutions,
and renders [GEMINI_FLASH_PROMPT: ...] tags as a highlighted placeholder line.
PPT suits shorter practice sets; for full 65-Q mocks .docx/.pdf read better.
Schema: references/paper_spec_schema.md.

Dependencies: python-pptx.
"""
from __future__ import annotations

import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

import paperlib as P

TEAL = RGBColor(0x0F, 0x8B, 0x8D)
NAVY = RGBColor(0x1F, 0x2A, 0x44)
AMBER = RGBColor(0xD9, 0x77, 0x06)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREY = RGBColor(0x6B, 0x72, 0x80)
INK = RGBColor(0x2A, 0x2F, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT_TEAL = RGBColor(0xE6, 0xF4, 0xF4)
TINT_AMBER = RGBColor(0xFF, 0xF4, 0xE0)

EMU = 914400
SW, SH = Inches(13.333), Inches(7.5)


def _bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _txt(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, shrink=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if shrink:
        # PowerPoint's "Shrink text on overflow" — keeps long content on the slide
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return tf


def _line(tf, runs, *, first=False, space_after=4, align=PP_ALIGN.LEFT, level=0):
    """runs = list of (text, {bold,color,size,italic})."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    for text, fmt in runs:
        r = p.add_run()
        r.text = text
        r.font.bold = fmt.get("bold", False)
        r.font.italic = fmt.get("italic", False)
        r.font.size = Pt(fmt.get("size", 16))
        r.font.color.rgb = fmt.get("color", INK)
        r.font.name = "Calibri"
    return p


def _clean_tag(tag):
    return P.gemini_prompt(tag)   # inner content + uniform style block


def _new_slide(prs, section_label, top_right=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _rect(slide, 0, 0, SW, Inches(0.62), NAVY)
    _rect(slide, 0, 0, Inches(0.18), SH, TEAL)
    tf = _txt(slide, Inches(0.4), Inches(0.06), Inches(9.5), Inches(0.5), MSO_ANCHOR.MIDDLE)
    _line(tf, [(section_label, {"bold": True, "color": WHITE, "size": 16})], first=True, space_after=0)
    if top_right:
        tr = _txt(slide, Inches(10.0), Inches(0.06), Inches(3.0), Inches(0.5), MSO_ANCHOR.MIDDLE)
        _line(tr, [(top_right, {"bold": True, "color": TINT_TEAL, "size": 13})],
              first=True, space_after=0, align=PP_ALIGN.RIGHT)
    return slide


def title_slide(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _rect(slide, 0, Inches(2.7), SW, Inches(0.06), TEAL)
    tf = _txt(slide, Inches(0.9), Inches(2.9), Inches(11.5), Inches(2.5))
    _line(tf, [("EXAMDOST  ·  GENERATED QUESTION BANK", {"bold": True, "color": TEAL, "size": 16})],
          first=True, space_after=10)
    _line(tf, [(P.render_math(meta.get("paper_title") or meta.get("exam", "Question Paper")),
                {"bold": True, "color": WHITE, "size": 34})], space_after=10)
    bits = []
    for k, suf in [("exam", ""), ("total_questions", " Questions"), ("total_marks", " Marks"),
                   ("duration_min", " min")]:
        if meta.get(k):
            bits.append(f"{meta[k]}{suf}")
    if meta.get("generated_on"):
        bits.append(f"Generated {meta['generated_on']}")
    _line(tf, [("    |    ".join(str(b) for b in bits), {"color": TINT_TEAL, "size": 14})])


def _meta_tagline(q, tags):
    if not tags:
        return ""
    label_map = {"Subject": "subject", "Chapter": "chapter", "Topic": "topic",
                 "Marks": "marks", "Type": "type"}
    parts = [f"{t}: {q.get(label_map[t])}" for t in tags
             if t in label_map and q.get(label_map[t]) not in (None, "")]
    return "    ".join(parts)


def add_stem(tf, q, tags, first=True):
    head = [(f"Q{q.get('number','')}.  ", {"bold": True, "color": TEAL, "size": 20})]
    if q.get("is_wildcard"):
        head.append(("[WILDCARD]  ", {"bold": True, "color": AMBER, "size": 14}))
    head.append((P.render_math(str(q.get("text", ""))), {"color": INK, "size": 18}))
    mk = q.get("marks")
    if mk not in (None, ""):
        head.append((f"   [{mk} mark{'s' if str(mk) != '1' else ''}]",
                     {"bold": True, "color": GREY, "size": 12}))
    _line(tf, head, first=first, space_after=8)
    for tag in q.get("diagram_prompts") or []:
        _line(tf, [("[DIAGRAM] ", {"bold": True, "color": AMBER, "size": 12}),
                   (_clean_tag(tag), {"italic": True, "color": INK, "size": 12})], space_after=6)
    for opt in q.get("options") or []:
        _line(tf, [(P.render_math(str(opt)), {"color": INK, "size": 15})], space_after=3, level=1)
    tagline = _meta_tagline(q, tags)
    if tagline:
        _line(tf, [(tagline, {"italic": True, "color": GREY, "size": 11})], space_after=4)


def add_answer_solution(tf, q, show_solutions, first=False):
    _line(tf, [("Answer Key:  ", {"bold": True, "color": GREEN, "size": 15}),
               (P.render_math(str(q.get("answer", ""))), {"bold": True, "color": INK, "size": 15})],
          first=first, space_after=4)
    if q.get("concept"):
        _line(tf, [("Concept:  ", {"bold": True, "color": TEAL, "size": 13}),
                   (P.render_math(str(q["concept"])), {"color": INK, "size": 13})], space_after=4)
    if not show_solutions:
        return
    sol = q.get("solution") or {}
    if not sol:
        return
    _line(tf, [("Solution", {"bold": True, "color": NAVY, "size": 13})], space_after=2)
    for label, key in [("Step 1 — Given", "given"), ("Step 2 — Formula", "formula"),
                       ("Step 3 — Calculation", "calculation")]:
        payload = sol.get(key)
        if payload in (None, "", []):
            continue
        joined = " ;  ".join(payload) if isinstance(payload, list) else str(payload)
        _line(tf, [(f"{label}: ", {"bold": True, "color": NAVY, "size": 12}),
                   (P.render_math(joined), {"color": INK, "size": 12})], space_after=2, level=1)
    if sol.get("final_answer"):
        _line(tf, [("Step 4 — Final Answer: ", {"bold": True, "color": GREEN, "size": 12}),
                   (P.render_math(str(sol["final_answer"])), {"bold": True, "color": INK, "size": 12})],
              space_after=2, level=1)
    for tag in sol.get("solution_diagram_prompts") or []:
        _line(tf, [("[DIAGRAM] ", {"bold": True, "color": AMBER, "size": 11}),
                   (_clean_tag(tag), {"italic": True, "color": INK, "size": 11})], space_after=2, level=1)


def _content_lines(q, with_solution):
    """Rough line estimate to decide whether a question overflows one slide."""
    lines = max(1, len(str(q.get("text", ""))) // 55)
    lines += len(q.get("options") or [])
    lines += 2 * len(q.get("diagram_prompts") or [])
    if with_solution:
        sol = q.get("solution") or {}
        if sol:
            lines += 2  # answer + concept
            for k in ("given", "formula", "calculation"):
                v = sol.get(k)
                lines += (len(v) if isinstance(v, list) else 1) if v else 0
            lines += 1 + len(sol.get("solution_diagram_prompts") or [])
    return lines


# slides taller than this (in estimated lines) are split across a continuation slide
_LINE_BUDGET = 18


def _content_box(slide):
    return _txt(slide, Inches(0.45), Inches(0.85), Inches(12.4), Inches(6.3), shrink=True)


def build(spec, out_path):
    out_path = P.safe_output_path(out_path)
    meta = spec.get("meta", {})
    layout = (meta.get("layout") or "combined").lower()
    show_solutions = meta.get("show_solutions", True)
    tags = meta.get("metadata_tags") or []
    sections = spec.get("sections", [])

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    title_slide(prs, meta)

    def solution_slide(sec_name, q, label="— Solution"):
        slide = _new_slide(prs, sec_name + f" {label}", f"Q{q.get('number','')}")
        tf = _content_box(slide)
        _line(tf, [(f"Q{q.get('number','')}.  ", {"bold": True, "color": TEAL, "size": 16}),
                   (P.render_math(str(q.get("text", ""))), {"color": INK, "size": 13})],
              first=True, space_after=6)
        add_answer_solution(tf, q, show_solutions, first=False)

    if layout == "split":
        for sec in sections:
            for q in sec.get("questions", []):
                slide = _new_slide(prs, sec.get("section", "Section"), f"Q{q.get('number','')}")
                add_stem(_content_box(slide), q, tags, first=True)
        for sec in sections:
            for q in sec.get("questions", []):
                solution_slide(sec.get("section", "Section"), q)
    else:
        for sec in sections:
            for q in sec.get("questions", []):
                # if a fully-combined slide would overflow, split: stem here, solution on a cont. slide
                if show_solutions and _content_lines(q, True) > _LINE_BUDGET:
                    slide = _new_slide(prs, sec.get("section", "Section"), f"Q{q.get('number','')}")
                    add_stem(_content_box(slide), q, tags, first=True)
                    solution_slide(sec.get("section", "Section"), q, label="(cont.)")
                else:
                    slide = _new_slide(prs, sec.get("section", "Section"), f"Q{q.get('number','')}")
                    tf = _content_box(slide)
                    add_stem(tf, q, tags, first=True)
                    add_answer_solution(tf, q, show_solutions, first=False)

    prs.save(out_path)
    nq = sum(len(s.get("questions", [])) for s in sections)
    print(f"Wrote {out_path}  ({nq} questions, layout={layout}, solutions={'on' if show_solutions else 'off'})")


def main():
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        spec = json.load(f)
    build(spec, sys.argv[2])


if __name__ == "__main__":
    main()
